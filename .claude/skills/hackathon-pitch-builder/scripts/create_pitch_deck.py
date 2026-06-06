#!/usr/bin/env python3
"""
AI BEAVERS Hackathon — 7-slide pitch deck generator.

Adapted from the investor create_pitch_deck.py. Three changes vs. the fork:

  1. Exactly the 7 work-slides the hackathon mandates (Problem+Customer,
     Solution+Product, Why Now, Market+Competition, Business+Evidence,
     Go-to-Market, Team). No separate title/cover slide — company name +
     tagline render as a banner on slide 1, so the deck stays at 7 and every
     slide does work ("max 7 slides").
  2. German slide titles (the judges + the live pitch are in German; the
     product canon is German).
  3. Slide 1 carries the one-line "first sentence" framing; slide 2 carries an
     optional demo URL / QR-line; Market+Competition is a combined two-column
     slide; Team slide ends with contact info.

Usage:
    python3 create_pitch_deck.py <data.json> [output.pptx]
"""

import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Palette: deep navy / blue / red accent (trust + clarity + urgency) -----
PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
SECONDARY = RGBColor(0x2B, 0x6C, 0xB0)
ACCENT = RGBColor(0xE5, 0x3E, 0x3E)
TEXT = RGBColor(0x1A, 0x20, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xFA, 0xFC)
DIVIDER = RGBColor(0xE2, 0xE8, 0xF0)

RECT, OVAL = 1, 3
MAX_SLIDES = 7


def create_pitch_deck(data, output_file="pitch_deck.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 — projector-friendly
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    company = data.get("company_name", "Projektname")
    W = 13.333

    state = {"n": 0}

    def _bg(slide, color):
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = color

    def _shape(slide, kind, x, y, w, h, color, line=None, line_w=1):
        sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = Pt(line_w)
        return sh

    def _footer(slide):
        state["n"] += 1
        box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(W - 1.2), Inches(0.35))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = company
        run.font.size = Pt(10)
        run.font.color.rgb = SECONDARY
        num = slide.shapes.add_textbox(Inches(W - 0.9), Inches(7.05), Inches(0.7), Inches(0.35))
        pn = num.text_frame.paragraphs[0]
        pn.alignment = PP_ALIGN.RIGHT
        rn = pn.add_run()
        rn.text = f"{state['n']}/7"
        rn.font.size = Pt(10)
        rn.font.color.rgb = SECONDARY

    def _header(slide, title, kicker=None):
        _bg(slide, WHITE)
        _shape(slide, RECT, 0, 0, W, 1.25, PRIMARY)
        _shape(slide, OVAL, W - 1.0, 0.25, 0.75, 0.75, ACCENT)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(W - 2.0), Inches(0.85))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.size = Pt(38)
        r.font.bold = True
        r.font.color.rgb = WHITE
        if kicker:
            kb = slide.shapes.add_textbox(Inches(0.52), Inches(0.92), Inches(W - 2.0), Inches(0.3))
            kr = kb.text_frame.paragraphs[0].add_run()
            kr.text = kicker
            kr.font.size = Pt(13)
            kr.font.color.rgb = DIVIDER

    def add_content_slide(title, items, kicker=None, lead=None):
        slide = prs.slides.add_slide(blank)
        _header(slide, title, kicker)
        _shape(slide, RECT, 0.5, 1.55, W - 1.0, 5.15, LIGHT, line=DIVIDER, line_w=1)
        top = 1.95
        if lead:
            lb = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(W - 1.8), Inches(0.9))
            lb.text_frame.word_wrap = True
            lr = lb.text_frame.paragraphs[0].add_run()
            lr.text = lead
            lr.font.size = Pt(22)
            lr.font.bold = True
            lr.font.italic = True
            lr.font.color.rgb = PRIMARY
            top += 1.1
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(W - 1.8), Inches(6.5 - top))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = f"•  {item}"
            r.font.size = Pt(20)
            r.font.color.rgb = TEXT
            p.space_after = Pt(14)
        _footer(slide)
        return slide

    def add_two_column_slide(title, lt, lc, rt, rc, kicker=None):
        slide = prs.slides.add_slide(blank)
        _header(slide, title, kicker)
        half = (W - 1.3) / 2
        for x, head, content, accent in (
            (0.5, lt, lc, SECONDARY),
            (0.8 + half, rt, rc, PRIMARY),
        ):
            _shape(slide, RECT, x, 1.55, half, 5.15, LIGHT, line=accent, line_w=2)
            tb = slide.shapes.add_textbox(Inches(x + 0.3), Inches(1.8), Inches(half - 0.6), Inches(4.7))
            tf = tb.text_frame
            tf.word_wrap = True
            hr = tf.paragraphs[0].add_run()
            hr.text = head
            hr.font.size = Pt(22)
            hr.font.bold = True
            hr.font.color.rgb = accent
            tf.paragraphs[0].space_after = Pt(10)
            for item in content:
                p = tf.add_paragraph()
                r = p.add_run()
                r.text = f"•  {item}"
                r.font.size = Pt(16)
                r.font.color.rgb = TEXT
                p.space_after = Pt(9)
        _footer(slide)
        return slide

    def _items(key):
        v = data.get(key)
        if v is None:
            return None
        return v if isinstance(v, list) else [v]

    # --- The 7 mandated work-slides ---------------------------------------
    # Slide 1 — Problem + Customer (carries company + first-sentence framing)
    kicker1 = f"{company}" + (f" — {data['tagline']}" if data.get("tagline") else "")
    add_content_slide(
        "Problem & Kunde",
        _items("problem_customer") or ["<Problem + namentlicher Kunde>"],
        kicker=kicker1,
        lead=data.get("first_sentence"),
    )

    # Slide 2 — Solution + Product (+ demo line)
    sol = _items("solution_product") or ["<Was ihr gebaut habt>"]
    if data.get("demo_url"):
        sol = sol + [f"Live-Demo: {data['demo_url']}"]
    add_content_slide("Lösung & Produkt", sol, kicker="Was wir gebaut haben — live zeigbar")

    # Slide 3 — Why Now
    add_content_slide("Why Now", _items("why_now") or ["<Was sich geändert hat>"],
                      kicker="Warum gerade jetzt möglich / dringend")

    # Slide 4 — Market + Competition (combined two-column)
    mc = data.get("market_competition", {})
    if isinstance(mc, dict):
        market = mc.get("market", ["<Bottom-up: Kundenzahl × €/Kunde>"])
        advantages = mc.get("our_advantages", ["<Warum wir gewinnen>"])
        competitors = mc.get("competitors", ["<Status quo / direkte + indirekte>"])
        add_two_column_slide(
            "Markt & Wettbewerb",
            "Markt (bottom-up)", market,
            "Wettbewerb & Vorsprung", competitors + ["—"] + advantages,
            kicker="Kundenzahl × Umsatz/Kunde — nie '1% von $100B'",
        )
    elif mc:
        add_content_slide("Markt & Wettbewerb", mc if isinstance(mc, list) else [mc])

    # Slide 5 — Business model + Traction / Evidence
    add_content_slide("Geschäftsmodell & Evidence",
                      _items("business_evidence") or ["<Wer zahlt, wie viel> + <was validiert>"],
                      kicker="Wer zahlt, wie oft — und was ihr validiert habt")

    # Slide 6 — Go-to-market
    add_content_slide("Go-to-Market",
                      _items("go_to_market") or ["<Wie die ersten 50 Kunden euch finden>"],
                      kicker="Ein Kanal, den ihr wirklich fahren könnt")

    # Slide 7 — Team (+ contact)
    team = _items("team") or ["<Warum dieses Team den Edge hat>"]
    if data.get("contact"):
        team = team + [f"Kontakt: {data['contact']}"]
    add_content_slide("Team", team, kicker="Warum dieses Team — Kontakt am Ende")

    if state["n"] > MAX_SLIDES:
        print(f"⚠️  WARNING: {state['n']} slides — hackathon cap is {MAX_SLIDES}.")
    prs.save(output_file)
    print(f"✅ Deck erstellt: {output_file} ({state['n']} Slides, Limit {MAX_SLIDES})")
    return output_file


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 create_pitch_deck.py <data.json> [output.pptx]")
        sys.exit(1)
    out = sys.argv[2] if len(sys.argv) > 2 else "pitch_deck.pptx"
    with open(sys.argv[1], encoding="utf-8") as f:
        create_pitch_deck(json.load(f), out)
